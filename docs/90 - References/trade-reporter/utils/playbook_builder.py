"""Template-first SMB PlayBook PPTX builder.

The builder keeps the official 17-slide SMB structure and hides unused slides
instead of deleting them. Text and images are written into the original template
zones so the disclosure/logo areas and slide design remain intact.
"""
from __future__ import annotations

import base64
import copy
import io
from datetime import datetime
from pathlib import Path
from typing import Iterable

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

TEMPLATE = Path(__file__).parent.parent / "assets" / "SMB_PlayBook_Template_2024.pptx"
EMU_IN = 914400


def _b64_stream(data_url: str | None) -> io.BytesIO | None:
    if not data_url:
        return None
    try:
        payload = data_url.split(",", 1)[-1]
        return io.BytesIO(base64.b64decode(payload))
    except (ValueError, TypeError, base64.binascii.Error):
        return None


def _iter_shapes(shapes) -> Iterable:
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _find_shape(slide, contains: str):
    needle = contains.casefold()
    for shape in _iter_shapes(slide.shapes):
        if shape.has_text_frame and needle in shape.text_frame.text.casefold():
            return shape
    return None


def _find_text_shapes(slide, contains: str) -> list:
    needle = contains.casefold()
    return [
        shape
        for shape in _iter_shapes(slide.shapes)
        if shape.has_text_frame and needle in shape.text_frame.text.casefold()
    ]


def _set_run_text(paragraph, text: str, *, size: int | None = None,
                  bold: bool | None = None, font_name: str | None = None) -> None:
    """Change text while retaining the paragraph and first run formatting."""
    runs = list(paragraph.runs)
    if runs:
        run = runs[0]
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        run = paragraph.add_run()
    run.text = text
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if font_name:
        run.font.name = font_name


def _set_plain_text(shape, text: str, *, size: int | None = None,
                    bold: bool | None = None, font_name: str = "Lato") -> None:
    """Replace a non-bulleted text box, preserving its first paragraph style."""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    lines = (text or "").splitlines() or [""]
    existing = list(tf.paragraphs)
    template_p = existing[0]

    for index, line in enumerate(lines):
        if index < len(existing):
            paragraph = existing[index]
        else:
            paragraph = tf.add_paragraph()
            # Clone paragraph-level properties (alignment, spacing, etc.).
            src_ppr = template_p._p.pPr
            if src_ppr is not None:
                paragraph._p.insert(0, copy.deepcopy(src_ppr))
        _set_run_text(paragraph, line, size=size, bold=bold, font_name=font_name)

    for paragraph in list(tf.paragraphs[len(lines):]):
        _set_run_text(paragraph, "", size=size, bold=bold, font_name=font_name)


def _clean_lines(value: str | None, limit: int) -> list[str]:
    lines = [line.strip().lstrip("•-–— ") for line in (value or "").splitlines()]
    return [line for line in lines if line][:limit]


def _fill_template_paragraphs(shape, lines: list[str], *, size: int | None = None,
                              font_name: str = "Lato") -> None:
    """Fill existing template paragraphs so bullets and spacing are preserved."""
    if not shape or not shape.has_text_frame:
        return
    paragraphs = list(shape.text_frame.paragraphs)
    for index, paragraph in enumerate(paragraphs):
        value = lines[index] if index < len(lines) else ""
        _set_run_text(paragraph, value, size=size, font_name=font_name)


SECTION_TITLES = (
    "THE BIG PICTURE",
    "INTRADAY FUNDAMENTALS",
    "TECHNICAL ANALYSIS",
    "TRADE STRATEGY",
    "TRADE MANAGEMENT",
    "RISK MANAGEMENT",
    "READING THE TAPE",
    "TECHNOLOGY",
    "TRADE REVIEW",
    "SCORE CARD",
)


def _style_section_title(slide) -> None:
    """Template rule: section titles in Lato BOLD CAPS."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip().upper()
        if not text.startswith(SECTION_TITLES):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = run.text.upper()
                run.font.name = "Lato"
                run.font.bold = True


def _hide_slide(slide) -> None:
    """PowerPoint equivalent of Skip/Hide Slide (keeps template structure)."""
    slide._element.set("show", "0")


def _show_slide(slide) -> None:
    slide._element.attrib.pop("show", None)


def _picture_placeholders(slide) -> list:
    """Return non-logo pictures ordered left-to-right, top-to-bottom."""
    pictures = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        and (shape.top / EMU_IN > 1.0 or shape.width / EMU_IN > 5.0)
    ]
    return sorted(pictures, key=lambda shape: (shape.top, shape.left))


def _remove_shape(shape) -> None:
    if shape is not None:
        shape._element.getparent().remove(shape._element)


def _add_image_in_box(slide, stream: io.BytesIO, left: int, top: int,
                      width: int, height: int):
    """Fit and center an image inside the exact template box."""
    stream.seek(0)
    with Image.open(stream) as image:
        image_width, image_height = image.size
    scale = min(width / image_width, height / image_height)
    out_width = int(image_width * scale)
    out_height = int(image_height * scale)
    out_left = int(left + (width - out_width) / 2)
    out_top = int(top + (height - out_height) / 2)
    stream.seek(0)
    return slide.shapes.add_picture(
        stream, out_left, out_top, width=out_width, height=out_height
    )


def _replace_picture(slide, placeholder, stream: io.BytesIO | None,
                     hyperlink: str | None = None):
    if placeholder is None:
        return None
    left, top, width, height = (
        placeholder.left, placeholder.top, placeholder.width, placeholder.height
    )
    _remove_shape(placeholder)
    if not stream:
        return None
    picture = _add_image_in_box(slide, stream, left, top, width, height)
    if hyperlink:
        try:
            picture.click_action.hyperlink.address = hyperlink
        except Exception:
            pass
    return picture


def _format_date(value: str | None) -> str:
    if not value:
        return ""
    try:
        return datetime.strptime(value, "%Y-%m-%d").strftime("%B %d, %Y")
    except ValueError:
        return value


def _split_title(value: str | None) -> tuple[str, str]:
    title = (value or "NAME OF THE TRADE").strip().upper()
    explicit = [line.strip() for line in title.splitlines() if line.strip()]
    if len(explicit) >= 2:
        return explicit[0], " ".join(explicit[1:])
    if len(title) <= 20:
        return title, ""
    words = title.split()
    best_index = min(
        range(1, len(words)),
        key=lambda index: abs(len(" ".join(words[:index])) - len(" ".join(words[index:])))
    )
    return " ".join(words[:best_index]), " ".join(words[best_index:])


def _fill_title_slide(slide, data: dict) -> None:
    shape = _find_shape(slide, "NAME OF THE")
    if not shape:
        return
    paragraphs = list(shape.text_frame.paragraphs)
    line_1, line_2 = _split_title(data.get("trade_name"))
    values = {
        0: line_1,
        1: line_2,
        2: "",
        3: (data.get("trader") or "").strip(),
        4: (data.get("ticker") or "").strip().upper(),
        5: _format_date(data.get("date")),
        6: "",
    }
    for index, paragraph in enumerate(paragraphs):
        _set_run_text(paragraph, values.get(index, ""), font_name="Lato")


def _fill_big_picture(slides: list, data: dict, hidden: set[int]) -> None:
    # Slide 4 — optional divider/overall market chart.
    divider = slides[3]
    divider_stream = _b64_stream(data.get("bp_chart"))
    divider_pictures = _picture_placeholders(divider)
    if divider_stream and divider_pictures:
        _replace_picture(divider, divider_pictures[0], divider_stream)
    else:
        hidden.add(3)

    spy_note = (data.get("spy_note") or "").strip()
    qqq_note = (data.get("qqq_note") or "").strip()
    spy_stream = _b64_stream(data.get("spy_chart"))
    qqq_stream = _b64_stream(data.get("qqq_chart"))

    notes_slide = slides[4]
    one_chart_slide = slides[5]
    two_chart_slide = slides[6]

    # Auto-select the SMB layout that best matches supplied content.
    chart_count = int(spy_stream is not None) + int(qqq_stream is not None)
    if chart_count == 2:
        # Slide 5 = SPY/QQQ context notes, Slide 7 = both charts.
        note_shapes = sorted(_find_text_shapes(notes_slide, "Example:"), key=lambda s: s.left)
        if len(note_shapes) >= 2 and (spy_note or qqq_note):
            _set_plain_text(note_shapes[0], spy_note, size=25)
            _set_plain_text(note_shapes[1], qqq_note, size=25)
        else:
            hidden.add(4)
        hidden.add(5)
        placeholders = _picture_placeholders(two_chart_slide)
        if len(placeholders) >= 2:
            _replace_picture(two_chart_slide, placeholders[0], spy_stream)
            _replace_picture(two_chart_slide, placeholders[1], qqq_stream)
    elif chart_count == 1:
        # Slide 6 = one note + one large chart.
        hidden.add(6)
        ticker = "$SPY" if spy_stream else "$QQQ"
        note = spy_note if spy_stream else qqq_note
        stream = spy_stream or qqq_stream
        label = _find_shape(one_chart_slide, "$SPY")
        note_shape = _find_shape(one_chart_slide, "Example:")
        if label:
            _set_plain_text(label, ticker, size=30, bold=True)
        if note_shape:
            _set_plain_text(note_shape, note, size=25)
        placeholders = _picture_placeholders(one_chart_slide)
        if placeholders:
            _replace_picture(one_chart_slide, placeholders[0], stream)

        # Keep the dual-note slide only when the other index has useful context.
        other_note = qqq_note if spy_stream else spy_note
        if other_note:
            note_shapes = sorted(_find_text_shapes(notes_slide, "Example:"), key=lambda s: s.left)
            if len(note_shapes) >= 2:
                _set_plain_text(note_shapes[0], "" if spy_stream else spy_note, size=25)
                _set_plain_text(note_shapes[1], qqq_note if spy_stream else "", size=25)
        else:
            hidden.add(4)
    else:
        hidden.update({5, 6})
        note_shapes = sorted(_find_text_shapes(notes_slide, "Example:"), key=lambda s: s.left)
        if len(note_shapes) >= 2 and (spy_note or qqq_note):
            _set_plain_text(note_shapes[0], spy_note, size=25)
            _set_plain_text(note_shapes[1], qqq_note, size=25)
        else:
            hidden.add(4)


def _fill_fundamentals(slide, data: dict, hidden: set[int]) -> None:
    narrative = (data.get("fundamentals") or "").strip()
    stats = data.get("fundamentals_stats") or []
    clean_stats = [
        {
            "name": str(item.get("name") or "").strip(),
            "value": str(item.get("value") or "").strip(),
        }
        for item in stats[:10]
        if isinstance(item, dict)
    ]
    has_stats = any(item["name"] or item["value"] for item in clean_stats)
    if not narrative and not has_stats:
        hidden.add(7)
        return

    narrative_shape = _find_shape(slide, "Example:")
    if narrative_shape:
        _set_plain_text(narrative_shape, narrative, size=27)

    stat_groups = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.GROUP]
    stat_groups.sort(key=lambda shape: (shape.top, shape.left))
    for index, group in enumerate(stat_groups[:10]):
        item = clean_stats[index] if index < len(clean_stats) else {"name": "", "value": ""}
        name_shape = next(
            (shape for shape in _iter_shapes(group.shapes)
             if shape.has_text_frame and "Stat Name" in shape.text_frame.text),
            None,
        )
        value_shape = next(
            (shape for shape in _iter_shapes(group.shapes)
             if shape.has_text_frame and "Stat #" in shape.text_frame.text),
            None,
        )
        if name_shape:
            _set_plain_text(name_shape, item["name"], size=15, bold=True)
        if value_shape:
            _set_plain_text(value_shape, item["value"], size=20, bold=True)


def _fill_chart_slide(slide, data_url: str | None, hidden: set[int], index: int) -> None:
    stream = _b64_stream(data_url)
    placeholders = _picture_placeholders(slide)
    if stream and placeholders:
        _replace_picture(slide, placeholders[0], stream)
    else:
        hidden.add(index)


def _fill_risk(slide, data: dict, hidden: set[int]) -> None:
    amount = (data.get("risk_amount") or "").strip()
    daily_pct = (data.get("daily_stop_pct") or "").strip()
    risk_details = " ".join(_clean_lines(data.get("risk_details"), 4))
    stop_details = " ".join(_clean_lines(data.get("stop_details"), 4))
    legacy = (data.get("risk") or "").strip()
    if legacy and not any((amount, daily_pct, risk_details, stop_details)):
        risk_details = legacy
    if not any((amount, daily_pct, risk_details, stop_details)):
        hidden.add(12)
        return

    left = _find_shape(slide, "How much will you risk")
    if left:
        paragraphs = list(left.text_frame.paragraphs)
        values = [
            "How much will you risk on the trade?",
            amount,
            "",
            "What percentage of your daily stop loss will you risk?",
            daily_pct,
        ]
        # The template box may have fewer paragraphs than answers.
        while len(paragraphs) < len(values):
            paragraphs.append(left.text_frame.add_paragraph())
        for idx, paragraph in enumerate(paragraphs):
            value = values[idx] if idx < len(values) else ""
            answer = idx in (1, 4)
            _set_run_text(
                paragraph, value, size=30 if not answer else 34,
                bold=True if answer else None, font_name="Lato"
            )

    right = _find_shape(slide, "Be specific")
    if right:
        paragraphs = list(right.text_frame.paragraphs)
        values = [
            f"RISK — {risk_details}" if risk_details else "",
            "",
            "",
            f"STOPS — {stop_details}" if stop_details else "",
        ]
        for idx, paragraph in enumerate(paragraphs):
            _set_run_text(
                paragraph, values[idx] if idx < len(values) else "",
                size=28, font_name="Lato"
            )


def _fill_media_and_bullets(slide, *, image_data: str | None, text: str | None,
                            max_lines: int, hidden: set[int], index: int,
                            hyperlink: str | None = None, text_anchor: str) -> None:
    stream = _b64_stream(image_data)
    lines = _clean_lines(text, max_lines)
    if not stream and not lines:
        hidden.add(index)
        return

    placeholders = _picture_placeholders(slide)
    if placeholders:
        _replace_picture(slide, placeholders[0], stream, hyperlink=hyperlink)
    shape = _find_shape(slide, text_anchor)
    if shape:
        _fill_template_paragraphs(shape, lines, size=27 if index == 14 else 30)


def _fill_review(slide, data: dict, hidden: set[int]) -> None:
    points = _clean_lines(data.get("review_points"), 3)
    detail = (data.get("review_detail") or "").strip()
    if not points and not detail:
        hidden.add(15)
        return
    points_shape = _find_shape(slide, "Add key point")
    detail_shape = _find_shape(slide, "detailed info")
    if points_shape:
        # The template inserts blank spacer paragraphs between the three bullets.
        paragraphs = list(points_shape.text_frame.paragraphs)
        mapping = {0: points[0] if len(points) > 0 else "",
                   2: points[1] if len(points) > 1 else "",
                   4: points[2] if len(points) > 2 else ""}
        for idx, paragraph in enumerate(paragraphs):
            _set_run_text(paragraph, mapping.get(idx, ""), size=28, font_name="Lato")
    if detail_shape:
        _set_plain_text(detail_shape, detail, size=26)


def _fill_scorecard(slide, grades: dict, hidden: set[int]) -> None:
    if not grades:
        hidden.add(16)
        return
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        for row in shape.table.rows:
            label = row.cells[0].text.strip()
            grade = str(grades.get(label, "")).strip()
            cell = row.cells[1]
            cell.text = grade
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Lato"
                    run.font.size = Pt(14)
                    run.font.bold = True


def build_playbook(data: dict) -> bytes:
    if not TEMPLATE.exists():
        raise FileNotFoundError(
            "SMB template missing. Place SMB_PlayBook_Template_2024.pptx in assets/."
        )

    presentation = Presentation(str(TEMPLATE))
    slides = list(presentation.slides)
    if len(slides) < 17:
        raise ValueError("Unexpected SMB template: 17 slides are required.")

    hidden: set[int] = {0}  # Notes slide is retained but hidden.

    _fill_title_slide(slides[2], data)
    _fill_big_picture(slides, data, hidden)
    _fill_fundamentals(slides[7], data, hidden)
    _fill_chart_slide(slides[8], data.get("ta_chart"), hidden, 8)
    _fill_chart_slide(slides[9], data.get("ta_chart2"), hidden, 9)
    if 9 not in hidden:  # drop the "(extra slide)" marker when actually used
        extra_title = _find_shape(slides[9], "extra slide")
        if extra_title:
            _set_plain_text(extra_title, "TECHNICAL ANALYSIS", bold=True)

    strategy_lines = _clean_lines(data.get("strategy"), 8)
    if strategy_lines:
        _fill_template_paragraphs(
            _find_shape(slides[10], "Add key point"), strategy_lines, size=27
        )
    else:
        hidden.add(10)

    _fill_chart_slide(slides[11], data.get("tm_chart"), hidden, 11)
    _fill_risk(slides[12], data, hidden)

    _fill_media_and_bullets(
        slides[13], image_data=data.get("tape_chart"), text=data.get("tape"),
        max_lines=4, hidden=hidden, index=13, hyperlink=data.get("tape_url"),
        text_anchor="key points for tape",
    )
    _fill_media_and_bullets(
        slides[14], image_data=data.get("technology_chart"), text=data.get("technology"),
        max_lines=6, hidden=hidden, index=14, text_anchor="Add key point",
    )
    _fill_review(slides[15], data, hidden)
    _fill_scorecard(slides[16], data.get("grades") or {}, hidden)

    for index, slide in enumerate(slides):
        if index in hidden:
            _hide_slide(slide)
        else:
            _show_slide(slide)
            if index >= 3:  # section slides only, not notes/disclaimer/title
                _style_section_title(slide)

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()
